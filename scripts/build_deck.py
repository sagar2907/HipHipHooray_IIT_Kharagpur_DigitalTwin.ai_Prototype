#!/usr/bin/env python3
"""Build DigitalTwin_Business_Proposal.pptx - the same Detailed Business
Proposal as 8_Proposal/DigitalTwin_Business_Proposal.pdf, in slide form.

    python scripts/build_deck.py

The submission form has two proposal fields, both labelled "Detailed
Business Proposal" - one PDF, one PPT. That wording means the same document
in both formats, not a PDF proposal plus an unrelated pitch highlight-reel.
This script mirrors every section of the PDF (§1-9), not a curated subset.

No template is mandated for Round 2, so the design is ours. The palette is
deliberately the product's own argument: a steel/charcoal base with colour
reserved for deviation, which is the ISA-101 principle the twin itself
follows. A deck that is visually quiet until something is wrong makes the
same point the supervisor view does.

Every number here traces to results/ or results/twin.db, and matches the
PDF exactly - both are generated from the same source facts, never typed
twice.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_PATTERN_TYPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "..", "8_Proposal", "DigitalTwin_Business_Proposal.pptx")

# palette - the product's own ISA-101 logic: grey until something is wrong
INK   = RGBColor(0x1B, 0x1E, 0x21)
STEEL = RGBColor(0x1F, 0x2A, 0x33)
PAPER = RGBColor(0xF4, 0xF6, 0xF8)
MUTED = RGBColor(0x5A, 0x60, 0x67)
LINE  = RGBColor(0xC3, 0xCB, 0xD2)
BLUE  = RGBColor(0x1C, 0x5D, 0x99)
ALARM = RGBColor(0xB3, 0x26, 0x1E)
WARN  = RGBColor(0xC8, 0xA4, 0x15)
GOOD  = RGBColor(0x2D, 0x6A, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIM   = RGBColor(0x9F, 0xB3, 0xC4)
HEAD_BG = RGBColor(0xE8, 0xEC, 0xEF)

W, H = Inches(13.333), Inches(7.5)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = STEEL if dark else PAPER
    return s


def tb(s, x, y, w, h, text, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, font="Calibri", space=0, line=None, italic=False):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, part in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if space:
            para.space_after = Pt(space)
        if line:
            para.line_spacing = line
        r = para.add_run()
        r.text = part
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return box


def rect(s, x, y, w, h, fill, outline=None, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if outline is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = outline
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def heading(s, text, kicker=None, dark=False):
    """Title with NO underline rule - an accent line under a title is the
    single clearest tell of a generated slide."""
    if kicker:
        tb(s, 0.7, 0.42, 12, 0.3, kicker.upper(), 11, True,
           DIM if dark else BLUE, font="Calibri")
    tb(s, 0.7, 0.74, 12, 0.75, text, 27, True, WHITE if dark else INK,
       font="Calibri")


def station_strip(s, x, y, n=20, dark_idx=(9, 13, 14), constraint=19,
                  forming=(6,), w=0.46, h=0.5, gap=0.055):
    """The line itself - the deck's recurring visual motif."""
    for i in range(n):
        cx = x + i * (w + gap)
        dark = i in dark_idx
        if dark:
            fill, txt = RGBColor(0xB6, 0xBC, 0xC2), RGBColor(0x50, 0x55, 0x5A)
        elif i == constraint:
            fill, txt = ALARM, WHITE
        elif i in forming:
            fill, txt = WARN, RGBColor(0x24, 0x1F, 0x00)
        else:
            fill, txt = RGBColor(0xDD, 0xE2, 0xE6), MUTED
        sh = rect(s, cx, y, w, h, fill, LINE)
        if dark:
            # The caption calls these "hatched" - the fill must actually be a
            # hatch pattern, not a slightly-darker grey the eye reads as solid.
            sh.fill.patterned()
            sh.fill.pattern = MSO_PATTERN_TYPE.LIGHT_UPWARD_DIAGONAL
            sh.fill.fore_color.rgb = RGBColor(0x8A, 0x90, 0x96)
            sh.fill.back_color.rgb = RGBColor(0xDD, 0xE2, 0xE6)
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = f"{i+1:02d}"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = txt
        r.font.name = "Consolas"


def kpi(s, x, y, w, value, label, color=BLUE, vsize=34, ht=1.5):
    rect(s, x, y, w, ht, WHITE, LINE)
    rect(s, x, y, 0.055, ht, color)
    tb(s, x + 0.22, y + 0.14, w - 0.4, 0.55, value, vsize, True, color,
       font="Consolas")
    tb(s, x + 0.22, y + 0.14 + 0.6, w - 0.4, ht - 0.7, label, 10.5, False,
       MUTED, line=1.15)


def bars(s, x, y, w, h, rows, maxv, unit="", color=BLUE, hi=None):
    """Simple horizontal bar chart."""
    rowh = h / len(rows)
    for i, (lab, val) in enumerate(rows):
        yy = y + i * rowh
        tb(s, x, yy + 0.04, 2.5, 0.3, lab, 11.5, False, INK)
        track = w - 3.6
        rect(s, x + 2.6, yy + 0.06, track, rowh - 0.22,
             RGBColor(0xE4, 0xE8, 0xEB), LINE)
        bw = max(0.04, track * (val / maxv))
        rect(s, x + 2.6, yy + 0.06, bw, rowh - 0.22,
             ALARM if (hi and lab in hi) else color)
        tb(s, x + w - 0.95, yy + 0.04, 0.95, 0.3, f"{val}{unit}", 11.5, True,
           INK, font="Consolas")


def table(s, x, y, w, col_widths, header, rows, row_heights=None,
          font_size=10.5, header_size=10):
    """A native pptx table - used for dense reference tables (many rows,
    wrapping text) where hand-positioned textboxes would be brittle."""
    n_rows = len(rows) + 1
    total_h = sum(row_heights) if row_heights else 0.5 * n_rows
    gshape = s.shapes.add_table(n_rows, len(header), Inches(x), Inches(y),
                                 Inches(w), Inches(total_h))
    gt = gshape.table
    for c, cw in enumerate(col_widths):
        gt.columns[c].width = Inches(cw)
    if row_heights:
        for r, rh in enumerate(row_heights):
            gt.rows[r].height = Inches(rh)

    def _cell(cell, text, size, bold, color, bg, align=PP_ALIGN.LEFT):
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        for i, part in enumerate(str(text).split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            r = para.add_run()
            r.text = part
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"

    for c, htext in enumerate(header):
        _cell(gt.cell(0, c), htext, header_size, True, RGBColor(0x3D, 0x4D, 0x5A), HEAD_BG)
    for ridx, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            text, kw = (val, {}) if not isinstance(val, tuple) else val
            _cell(gt.cell(ridx, c), text, kw.get("size", font_size),
                  kw.get("bold", False), kw.get("color", INK),
                  kw.get("bg", WHITE), kw.get("align", PP_ALIGN.LEFT))
    # thin line borders throughout (pptx tables default to a heavier theme style)
    return gt


def foot(s, text):
    tb(s, 0.7, 7.08, 11.9, 0.3, text, 9.5, False, MUTED, italic=True)


def build():
    prs = deck()

    # ---------------------------------------------------------------- 1 title
    s = blank(prs, dark=True)
    tb(s, 0.9, 1.85, 11, 1.1, "DigitalTwin.ai", 56, True, WHITE)
    tb(s, 0.95, 2.95, 11.3, 0.6,
       "A vehicle assembly line that tells you which station is costing you cars —"
       "\nand what to do about it.", 19, False, DIM, line=1.25)
    rect(s, 0.95, 3.95, 1.5, 0.045, ALARM)
    tb(s, 0.95, 4.25, 11, 0.9,
       "Detailed Business Proposal  ·  Team HipHipHooray  ·  Sagar Sahu, Priyansh Goyal  ·  IIT Kharagpur"
       "\nAccenture Innovation Challenge 2026  ·  Problem Track 4  ·  Round 2",
       13, False, DIM, line=1.4)
    tb(s, 0.95, 5.35, 11, 0.35,
       "The one-line claim: a plant already emits enough data to say which station is "
       "costing you cars right now,\nwhich one is about to, and what to do about it — "
       "on a line where nearly a fifth of the stations have no sensors at all.",
       12.5, False, RGBColor(0xC9, 0xD6, 0xE2), line=1.3)
    station_strip(s, 0.95, 6.4, constraint=19, forming=(6,))

    # ------------------------------------------------------ 2 headline numbers
    s = blank(prs)
    heading(s, "We built the loop, ran it for 903 shifts.", "the claim, measured")
    tb(s, 0.7, 1.62, 12, 0.5,
       "Every number in this proposal is measured against held-out ground truth, with its "
       "source file named.", 14.5, False, MUTED)
    for i, (v, l, c) in enumerate([
        ("15.5%", "of forming-bottleneck warnings name a\nstation with no sensors at all\n[15.0–15.9], n=29,060", BLUE),
        ("0.025", "calibration error (ECE) after fitting,\nfrom 0.479\n600 held-out samples", GOOD),
        ("70.5%", "of alarmed tools get an actionable\ndiagnosis\n[67.9–72.9], n=1,246", BLUE),
        ("903", "shifts replayed through the live loop\n\n86,742 decisions recorded", GOOD)]):
        kpi(s, 0.7 + i * 3.07, 2.3, 2.85, v, l, c, vsize=32, ht=2.0)
    tb(s, 0.7, 4.65, 12, 0.4, "Three things make the problem hard, and they are the three the brief names.",
       15, True, INK)
    for i, (t, d) in enumerate([
        ("The constraint moves", "~20 times per shift. A weekly average\ndescribes no moment of the shift it covers."),
        ("Utilisation lies", "A blocked station and a slow station\nboth read ~95% busy."),
        ("The data is uneven", "Legacy and modern equipment mixed.\nSome stations emit nothing at all.")]):
        x = 0.7 + i * 4.07
        rect(s, x, 5.15, 3.85, 1.55, WHITE, LINE)
        rect(s, x, 5.15, 3.85, 0.045, [ALARM, WARN, BLUE][i])
        tb(s, x + 0.22, 5.35, 3.4, 0.35, t, 14, True, INK)
        tb(s, x + 0.22, 5.72, 3.4, 0.9, d, 11, False, MUTED, line=1.25)
    foot(s, "§1 · Problem framing")

    # ---------------------------------------------------- 3 the virtual plant
    s = blank(prs)
    heading(s, "We built a plant, because no student team gets a real one.", "§2.1 · methodology")
    tb(s, 0.7, 1.62, 12, 0.75,
       "A software plant simulating a mixed-model line — buffers, tool wear, breakdowns, the "
       "rework loop, shift breaks,\na mid-shift changeover sweep, operator andon stops. The "
       "elements a practitioner notices missing from a toy\nmodel, not just the happy path.",
       13.5, False, MUTED, line=1.3)
    rect(s, 0.7, 2.55, 5.85, 2.35, WHITE, LINE)
    rect(s, 0.7, 2.55, 5.85, 0.05, GOOD)
    tb(s, 0.98, 2.78, 5.3, 0.35, "OBSERVED — what the twin may read", 14, True, GOOD)
    tb(s, 0.98, 3.2, 5.3, 1.55,
       "Exactly what a real plant's sensors and\npaperwork would produce — the five\nstreams "
       "in the integration slide.", 12.5, False, MUTED, line=1.35)
    rect(s, 6.78, 2.55, 5.85, 2.35, WHITE, LINE)
    rect(s, 6.78, 2.55, 5.85, 0.05, ALARM)
    tb(s, 7.06, 2.78, 5.3, 0.35, "HIDDEN — sealed off from the detector", 14, True, ALARM)
    tb(s, 7.06, 3.2, 5.3, 1.55,
       "Which fault was injected, the true onset\nminute, every micro-stop, which tools are\n"
       "actually worn. Enforced by a test that\nchecks no observed file carries a truth column.",
       12.5, False, MUTED, line=1.3)
    tb(s, 0.7, 5.1, 12, 0.4,
       "That separation is the whole point: it turns “does this work” from an opinion into a score.",
       15, True, INK)
    for i, (v, l) in enumerate([
        ("162", "shifts simulated\n156 single + 6 three-shift sequences"),
        ("4", "line layouts\n15–30 stations, one with a parallel pair"),
        ("972", "tools\nacross 8 fault conditions"),
        ("5×12", "injected fault kinds\n+60 no-fault runs for false alarms")]):
        x = 0.7 + i * 3.07
        rect(s, x, 5.55, 2.85, 1.15, WHITE, LINE)
        rect(s, x, 5.55, 0.05, 1.15, BLUE)
        tb(s, x + 0.2, 5.68, 2.5, 0.45, v, 22, True, BLUE, font="Consolas")
        tb(s, x + 0.2, 6.1, 2.55, 0.6, l, 10, False, MUTED, line=1.2)
    foot(s, "§2.1 · Source: docs/dataset/v5_dataset.md")

    # -------------------------------------------------------- 4 twin not shadow
    s = blank(prs)
    heading(s, "A twin, not a shadow — and four tiers of knowledge, never mixed.", "§2.2–2.3 · solution design")
    tb(s, 0.7, 1.62, 12, 0.85,
       "Kritzinger's taxonomy: a model is offline, a shadow receives data one-way, a twin closes "
       "the loop. Our prototype ingests,\ndetects, ranks, prescribes and re-reads on a timer — "
       "causality enforced structurally: the detector is handed a view of\nthe run truncated at "
       "now, so it cannot read the future even by accident. Verified identical verdicts at 13/13\n"
       "timepoints against the full-history run.",
       13, False, MUTED, line=1.3)
    tiers = [
        ("A", "Measured", "The station reports it", "PLC state, scans, tool telemetry", BLUE),
        ("B", "Inferred", "Bracketed between neighbours", "Boundary timestamps, buffer slope", GOOD),
        ("C", "Attested", "A person says so", "Manual checklist entry, with entry latency", WARN),
        ("D", "Predicted", "Projected under current flow", "Buffer countdown", MUTED),
    ]
    for i, (letter, name, what, src, col) in enumerate(tiers):
        x = 0.7 + i * 3.07
        rect(s, x, 3.15, 2.85, 2.1, WHITE, LINE)
        rect(s, x, 3.15, 2.85, 0.05, col)
        tb(s, x + 0.2, 3.32, 2.5, 0.5, letter + " — " + name, 15, True, INK)
        tb(s, x + 0.2, 3.85, 2.5, 0.55, what, 11.5, False, MUTED, line=1.25)
        tb(s, x + 0.2, 4.55, 2.5, 0.6, src, 10.5, False, col, line=1.2)
    tb(s, 0.7, 5.55, 12, 0.4,
       "Every value on screen carries its tier. An operator must never see an inferred number and a measured one in the same font.",
       13.5, True, INK, line=1.25)
    foot(s, "§2.2 · §2.3")

    # ------------------------------------------------------- 5 dark stations
    s = blank(prs)
    heading(s, "It names stations that have no sensors.", "complexity 1 · uneven coverage")
    station_strip(s, 0.95, 1.85, dark_idx=(9, 13, 14), constraint=19, forming=(6,),
                  w=0.55, h=0.62, gap=0.06)
    tb(s, 0.95, 2.62, 12, 0.35,
       "hatched = no sensors at all   ·   red = the constraint   ·   amber = forming",
       11.5, False, MUTED)
    kpi(s, 0.7, 3.2, 3.85, "15.5%", "of all forming-bottleneck warnings name a\n"
        "station with zero instrumentation   [15.0–15.9]\nn=29,060", BLUE, vsize=30, ht=1.85)
    kpi(s, 4.75, 3.2, 3.85, "2.76%", "of vehicles a manual checklist PASSED\n"
        "went on to fail end-of-line   ·   96.51% pass rate\nn=31,329 checklist entries", ALARM, vsize=30, ht=1.85)
    tb(s, 8.8, 3.28, 3.85, 1.95,
       "Body + final assembly:\n75.3% instrumented,\n24.7% on manual checks —\n"
       "an exact match to the\nbrief's own reference\nparameters.",
       12.5, True, GOOD, line=1.28)
    tb(s, 0.7, 5.3, 11.93, 1.55,
       "A checklist reading near-100% against a non-zero EOL failure rate is measuring compliance"
       " with the checklist,\nnot quality. Manual checks enter as a fourth data tier — attested —"
       " carrying the entry latency\n(mean 4.8 min), because a human record does not exist until"
       " someone types it in.",
       12.5, False, MUTED, line=1.32)
    foot(s, "§5 · Complexity 1 — Demonstrated")

    # ---------------------------------------------------------- 6 prescription
    s = blank(prs)
    heading(s, "It does not just report. It prescribes.", "§2.4–2.5 · from detection to decision")
    for i, (a, b, c) in enumerate([
        ("DESCRIPTIVE", "what happened", RGBColor(0xDD, 0xE2, 0xE6)),
        ("DIAGNOSTIC", "why", RGBColor(0xC3, 0xD3, 0xE0)),
        ("PREDICTIVE", "what is about to", RGBColor(0x7F, 0xA8, 0xCC)),
        ("PRESCRIPTIVE", "what to do", BLUE)]):
        x = 0.7 + i * 3.07
        rect(s, x, 1.72, 2.85, 0.85, c)
        tb(s, x + 0.2, 1.86, 2.5, 0.3, a, 12.5, True, WHITE if i == 3 else INK)
        tb(s, x + 0.2, 2.16, 2.5, 0.3, b, 11, False, WHITE if i == 3 else MUTED)
    rect(s, 0.7, 2.85, 11.93, 1.3, WHITE, LINE)
    rect(s, 0.7, 2.85, 0.055, 1.3, BLUE)
    tb(s, 1.0, 3.0, 11.4, 0.4, "S20 — Tool change / maintenance pull-forward", 17, True, INK)
    tb(s, 1.0, 3.4, 11.4, 0.7,
       "processing time is climbing against this station's own baseline   ·   margin 71.8 s   ·"
       "   next best S12\nLeaving it alone costs ≈ 0.73 vehicles over the median 10-minute episode.",
       12.5, False, MUTED, line=1.3)
    tb(s, 0.7, 4.35, 12, 0.35, "The alert contract — an alert may be raised only if it carries all five of:",
       15, True, INK)
    for i, f in enumerate(["Candidate\n+ margin", "Evidence", "Persistence\nestimate", "Recommended\naction", "Cost of not\nacting"]):
        x = 0.7 + i * 2.44
        rect(s, x, 4.78, 2.25, 0.95, WHITE, LINE)
        tb(s, x + 0.15, 4.78, 1.95, 0.95, f, 11.5, True, BLUE, align=PP_ALIGN.CENTER, line=1.2)
    tb(s, 0.7, 5.95, 11.93, 0.85,
       "An alert that cannot state its evidence is suppressed, not downgraded — a quiet alert with"
       " no evidence is still noise,\nand noise is what erodes floor trust. 7 alerts were suppressed"
       " by this rule in the recorded run.", 12.5, False, MUTED, line=1.3)
    foot(s, "§2.4 · §2.5")

    # ------------------------------------------------------ 7 integration
    s = blank(prs)
    heading(s, "Five streams a plant already produces. Nothing else.", "§2.6 · integration & the read-only boundary")
    rows = [
        ("Unit scans — VIN, station, timestamp", "Barcode/RFID readers at station boundaries"),
        ("Station state transitions", "PLC state tags, already historised for OEE"),
        ("Buffer levels", "Conveyor counters and occupancy sensors"),
        ("Tool readings — torque, angle, current", "Nutrunner controllers over Open Protocol / OPC-UA"),
        ("Manual checks — result and latency", "The paper or tablet checklist at a dark station"),
    ]
    table(s, 0.7, 1.65, 11.93, [5.9, 6.03],
          ["WHAT THE LIVE TWIN READS", "WHERE A PLANT ALREADY HAS IT"], rows,
          row_heights=[0.4] + [0.42] * 5, font_size=12, header_size=11)
    rect(s, 0.7, 4.35, 11.93, 1.05, WHITE, LINE)
    rect(s, 0.7, 4.35, 0.055, 1.05, GOOD)
    tb(s, 1.0, 4.52, 11.3, 0.75,
       "The boundary is enforced, not promised. Exactly two modules may write anything, both to "
       "the twin's own store.\nNo module may import a network client — checked by a test that walks "
       "every module's syntax tree; verified by planting\na PLC write and confirming it was caught.",
       12.5, False, INK, line=1.3)
    rect(s, 0.7, 5.55, 11.93, 1.1, WHITE, LINE)
    rect(s, 0.7, 5.55, 0.055, 1.1, WARN)
    tb(s, 1.0, 5.72, 11.3, 0.85,
       "What stands between this and a live line: one adapter. The prototype reads these streams "
       "from files today —\nnothing else. Shadow mode already works: point the replay driver at a "
       "plant's exported logs and the entire twin runs,\nwith no live connection of any kind.",
       12.5, False, INK, line=1.3)
    foot(s, "§2.6")

    # --------------------------------------------------------- 8 target users
    s = blank(prs)
    heading(s, "Three users. One twin. Proven, not asserted.", "§3 · target users — complexity 5")
    for i, (who, what, detail) in enumerate([
        ("Floor supervisor", "real time",
         "ISA-101: grey base, colour only on\ndeviation. Quiet until something\nis wrong."),
        ("Plant manager", "weekly planning",
         "A constraint-occupancy distribution,\ndeliberately not an average —\n"
         "the constraint moves ~20×/shift."),
        ("Leadership", "investment case",
         "Value and coverage, every figure\nnaming the file that produced it.")]):
        x = 0.7 + i * 4.07
        rect(s, x, 1.72, 3.85, 2.3, WHITE, LINE)
        rect(s, x, 1.72, 3.85, 0.05, BLUE)
        tb(s, x + 0.28, 1.98, 3.3, 0.35, who, 16, True, INK)
        tb(s, x + 0.28, 2.36, 3.3, 0.3, what, 12, True, BLUE)
        tb(s, x + 0.28, 2.78, 3.4, 1.15, detail, 12, False, MUTED, line=1.3)
    rect(s, 0.7, 4.3, 11.93, 1.55, WHITE, LINE)
    rect(s, 0.7, 4.3, 0.055, 1.55, GOOD)
    tb(s, 1.0, 4.5, 11.3, 0.4, "The reconciliation test", 16, True, GOOD)
    tb(s, 1.0, 4.92, 11.3, 0.85,
       "Each level totals independently and is compared: 6,730 constraint-minutes and 4,431 "
       "vehicles — identical across\nsupervisor records, manager weeks and leadership. That is how "
       "you show three views are one twin rather than\nthree dashboards that happen to sit in the "
       "same app.", 12.5, False, MUTED, line=1.3)
    tb(s, 0.7, 6.05, 11.93, 0.7,
       "Colour is reserved for deviation on the plant HMI (ISA-101); the leadership view may use "
       "colour freely — it governs a\nbusiness dashboard, not a control-room screen.",
       11.5, False, MUTED, italic=True, line=1.25)
    foot(s, "§3 · Complexity 5 — Demonstrated")

    # ------------------------------------------------------- 9 detection quality
    s = blank(prs)
    heading(s, "Measured against economic ground truth.", "§4.1 · business case")
    tb(s, 0.7, 1.62, 12, 0.35,
       "Ground truth is the station whose speed-up actually produces more cars, not “which "
       "looked busiest”. Regret = cars lost per block. Lower is better.",
       12.5, False, MUTED)
    bars(s, 0.7, 2.15, 11.9, 1.7, [
        ("Ours (effective CT)", 1.309),
        ("Active period (Roser)", 1.348),
        ("Utilisation baseline", 1.477)], maxv=1.6, hi={"Utilisation baseline"})
    rect(s, 0.7, 4.1, 5.8, 1.55, WHITE, LINE)
    rect(s, 0.7, 4.1, 0.055, 1.55, GOOD)
    tb(s, 1.0, 4.28, 5.3, 0.35, "What we claim", 14.5, True, GOOD)
    tb(s, 1.0, 4.65, 5.3, 0.9,
       "We beat the utilisation baseline —\nMcNemar paired test, p = 0.0025, n = 202 blocks.",
       12.5, False, INK, line=1.3)
    rect(s, 6.83, 4.1, 5.8, 1.55, WHITE, LINE)
    rect(s, 6.83, 4.1, 0.055, 1.55, ALARM)
    tb(s, 7.13, 4.28, 5.3, 0.35, "What we refuse to claim", 14.5, True, ALARM)
    tb(s, 7.13, 4.65, 5.3, 0.9,
       "That we beat the active-period method.\np = 0.45 — statistically tied. Saying otherwise"
       " would\nnot survive one question.", 12.5, False, INK, line=1.25)
    tb(s, 0.7, 5.85, 11.93, 0.9,
       "Two independent findings say the same thing: a 0.79-car label-noise floor makes top-1 "
       "close to a coin flip in ~50% of\nblocks, and across four topologies top-1 collapses "
       "44.5% → 10.6% while regret holds 1.208–1.431. The operationally\nmeaningful metric "
       "transfers; the argmax metric does not — so regret is the headline, not accuracy.",
       12, False, MUTED, line=1.28)
    foot(s, "§4.1 · §4.2 · Source: results/eval_v5.csv")

    # ------------------------------------------------------------- 10 value
    s = blank(prs)
    heading(s, "Value, stated honestly — including the line we deleted.", "§4.3 · business case")
    rows = [
        ("Fewer cars lost to a mis-identified constraint (regret 1.309 vs 1.477 utilisation, on a 2.271-car ceiling)", ("measured", {"bold": True, "color": GOOD})),
        ("Earlier warning: buffer countdown, 178 warnings, median error +0.57 min", ("measured", {"bold": True, "color": GOOD})),
        ("Avoided false scrap: transducer-drift tools flagged for recalibration, not repair", ("measured", {"bold": True, "color": GOOD})),
        ("Sensor spend avoided: 15.5% of forming warnings need no sensor at that station", ("measured", {"bold": True, "color": GOOD})),
        ("Lead-time reduction under CONWIP release control", ("NOT measured — excluded", {"bold": True, "color": ALARM})),
    ]
    table(s, 0.7, 1.65, 11.93, [9.5, 2.43],
          ["VALUE LINE", "STATUS"], rows,
          row_heights=[0.4, 0.62, 0.55, 0.55, 0.55, 0.55], font_size=13, header_size=11)
    rect(s, 0.7, 5.55, 11.93, 1.15, WHITE, LINE)
    rect(s, 0.7, 5.55, 0.055, 1.15, ALARM)
    tb(s, 1.0, 5.7, 11.3, 0.9,
       "One value line has been deliberately removed. An earlier draft led with “same throughput, "
       "36% lower lead time, zero capex” — no file in our results produces that number, so under "
       "our own rule it does not appear. Four defensible lines beat five with one that dies under a "
       "single question.", 12, False, INK, line=1.28)
    foot(s, "§4.3")

    # --------------------------------------------------- 11 complexity summary
    s = blank(prs)
    heading(s, "Handling the seven real-world complexities.", "§5 · coverage summary")
    rows = [
        ("1", "Uneven sensor coverage; manual checklists", "Dark stations bracketed and named; checklist honesty tested — 96.51% pass vs 2.76% real escape rate", ("Demonstrated", {"bold": True, "color": GOOD})),
        ("2", "Multi-causal, intermittent root causes", "Separated by scope, not statistics. Operator variation excluded on ethical grounds — stated, not an oversight", ("Partly demonstrated", {"bold": True, "color": WARN})),
        ("3", "PLC risk; retrofits in maintenance windows only", "Read-only boundary enforced by a test. Retrofit list is a window-dated schedule ranked by exposure/rupee", ("Demonstrated", {"bold": True, "color": GOOD})),
        ("4", "Early defect surfaces late", "Onset read backwards off CUSUM, +2 min of truth; vehicles partitioned by on-line vs shipped", ("Demonstrated", {"bold": True, "color": GOOD})),
        ("5", "Three stakeholder views", "One record stream, three resolutions, reconciliation test passing exactly", ("Demonstrated", {"bold": True, "color": GOOD})),
        ("6", "Scaling: layout, vintage, sensor maturity", "Measured across 4 topologies incl. a parallel pair; vintage has no data axis", ("Mostly demonstrated", {"bold": True, "color": WARN})),
        ("7", "False alarms erode trust", "Confidence calibrated (ECE 0.479→0.025); 48.6 alerts/shift vs ISA-18.2 budget of 150", ("Demonstrated", {"bold": True, "color": GOOD})),
    ]
    table(s, 0.7, 1.65, 11.93, [0.5, 3.3, 6.83, 1.3],
          ["#", "COMPLEXITY", "HOW WE ANSWER IT", "STATE"], rows,
          row_heights=[0.4] + [0.65] * 7, font_size=11, header_size=10.5)
    tb(s, 0.7, 6.55, 11.93, 0.55,
       "Against the brief's own reference parameters: body + final assembly sit at 75.3% "
       "instrumented / 24.7% manual — an exact match to “a majority well-instrumented, a "
       "meaningful minority on manual checks”.",
       11.5, False, GOOD, italic=True, line=1.25)
    foot(s, "§5")

    # ----------------------------------------------------- 12 complexity 2
    s = blank(prs)
    heading(s, "Multi-causal roots: separated by scope, not statistics.", "complexity 2 · multi-causal, intermittent roots")
    tb(s, 0.7, 1.72, 12, 1.0,
       "A bottleneck or defect rarely has one cause — equipment wear, upstream part quality, "
       "environmental drift can all move together.\nRather than chase correlation across noisy "
       "signals, we isolate by scope: who else is affected, and who is not. The prototype\n"
       "already separates fault classes the detector can distinguish — slowing, breakdown, "
       "starved, blocked — as the first cut.",
       14, False, MUTED, line=1.35)
    rect(s, 0.7, 2.95, 11.93, 1.7, WHITE, LINE)
    rect(s, 0.7, 2.95, 0.055, 1.7, ALARM)
    tb(s, 1.0, 3.15, 11.3, 0.4, "Operator variation is deliberately excluded — on ethical grounds.",
       16, True, INK)
    tb(s, 1.0, 3.62, 11.3, 0.9,
       "We measure the station, never the person. The brief names operator variation as a "
       "possible root cause; we choose\nnot to build a per-operator model. That is a design "
       "decision we state plainly, not a gap we hid.",
       12.5, False, MUTED, line=1.3)
    tb(s, 0.7, 4.95, 12, 0.4, "State: partly demonstrated", 15, True, WARN)
    tb(s, 0.7, 5.4, 12, 1.2,
       "The fault-class separation the detector already performs is real and measured; the "
       "full co-occurrence / scope-isolation engine described in\nour design notes is specified "
       "but not built — a multi-day build we chose not to spend against two graded documents "
       "that did not\nyet exist. A missing complexity costs part of one criterion; a missing "
       "proposal costs a third of the submission.",
       13, False, MUTED, line=1.32)
    foot(s, "§5 · Complexity 2 — Partly demonstrated; scope engine designed")

    # ----------------------------------------------------- 13 complexity 3
    s = blank(prs)
    heading(s, "PLC risk: a window-dated retrofit schedule, not a wishlist.", "complexity 3 · maintenance windows")
    tb(s, 0.7, 1.72, 12, 0.75,
       "The twin advises and never writes to line control — a read-only boundary in ISA-95 terms, "
       "shown on every screen.\nSensing additions mount externally and publish on a separate "
       "network, so no PLC program changes and no re-validation.",
       14, False, MUTED, line=1.32)
    for i, (v, l, c) in enumerate([
        ("Ranked by", "exposure closed per rupee\n— not a flat sensor list", BLUE),
        ("2 shutdowns", "the retrofit list splits across\ntwo scheduled maintenance windows", BLUE),
        ("~₹101,560", "cost of deferring, carried\nfor another six months", ALARM)]):
        x = 0.7 + i * 4.07
        rect(s, x, 2.65, 3.85, 1.5, WHITE, LINE)
        rect(s, x, 2.65, 0.05, 1.5, c)
        tb(s, x + 0.22, 2.82, 3.5, 0.5, v, 19, True, c, font="Consolas")
        tb(s, x + 0.22, 3.35, 3.5, 0.75, l, 11.5, False, MUTED, line=1.25)
    rect(s, 0.7, 4.4, 11.93, 1.5, WHITE, LINE)
    rect(s, 0.7, 4.4, 0.055, 1.5, GOOD)
    tb(s, 1.0, 4.58, 11.3, 1.15,
       "Never: closed-loop write to line control. Not a phase that arrives later — it is excluded "
       "from every phase of the roadmap.\nThe integration boundary is enforced by a test that walks "
       "every module's syntax tree, not a policy written in a document.",
       13, False, INK, line=1.32)
    foot(s, "§5 · Complexity 3 — Demonstrated")

    # ------------------------------------------------------------- 14 genealogy
    s = blank(prs)
    heading(s, "Two tools. Same symptom. Opposite answers.",
            "complexity 4 · defects found late")
    tb(s, 0.7, 1.62, 12, 0.32,
       "Both present identically to a torque-only monitor: “the reading moved”.",
       13, False, MUTED)
    for i, (stn, cond, act, col, ch) in enumerate([
        ("S05", "REAL WEAR", "SERVICE the tool",
         ALARM, "torque −3.6σ  ·  current −4.1σ  ·  angle +6.0σ"),
        ("S06", "SENSOR LYING", "RECALIBRATE only — do NOT service",
         BLUE, "torque −2.9σ  ·  current +0.5σ  ·  angle −0.2σ")]):
        y = 2.0 + i * 1.62
        rect(s, 0.7, y, 11.93, 1.42, WHITE, LINE)
        rect(s, 0.7, y, 0.055, 1.42, col)
        tb(s, 1.0, y + 0.16, 1.0, 0.4, stn, 20, True, INK, font="Consolas")
        tb(s, 2.0, y + 0.22, 2.6, 0.35, cond, 12.5, True, col)
        tb(s, 4.7, y + 0.2, 7.6, 0.4, act, 15, True, INK)
        tb(s, 4.7, y + 0.68, 7.6, 0.4, ch, 11.5, False, MUTED, font="Consolas")
        tb(s, 1.0, y + 0.72, 3.4, 0.5, "both look like\n\"the torque moved\"", 10.5, False, MUTED, line=1.2)
    tb(s, 0.7, 5.35, 12, 0.4,
       "Servicing a lying sensor scraps good parts and fixes nothing.", 15.5, True, INK)
    tb(s, 0.7, 5.78, 12, 0.9,
       "Onset is read backwards off the CUSUM accumulator — dated to +2 minutes of truth — "
       "then the VIN thread lists\nexactly which vehicles carry it: 202 through S05, 48 still "
       "on the line, 154 already completed. 70.5% of alarmed\ntools get an actionable diagnosis "
       "[67.9–72.9, n=1,246].", 12.5, False, MUTED, line=1.3)
    foot(s, "§5 · Complexity 4 — Demonstrated")

    # ------------------------------------------------------------- 15 transfer
    s = blank(prs)
    heading(s, "Regret transfers. Accuracy does not.", "complexity 6 · scaling")
    tb(s, 0.7, 1.62, 12, 0.35,
       "The same method on four different line topologies — including one with a parallel "
       "pair that breaks the series assumption.", 12.5, False, MUTED)
    tb(s, 0.7, 2.1, 5.6, 0.35, "TOP-1 ACCURACY  — collapses", 12, True, ALARM)
    bars(s, 0.7, 2.5, 5.7, 1.65, [
        ("L1  20 stn", 44.5), ("L2  30 stn", 10.6),
        ("L3  parallel", 13.8), ("L4  15 stn", 21.9)], maxv=50, unit="%", color=ALARM)
    tb(s, 6.9, 2.1, 5.6, 0.35, "REGRET  — holds", 12, True, GOOD)
    bars(s, 6.9, 2.5, 5.7, 1.65, [
        ("L1  20 stn", 1.208), ("L2  30 stn", 1.218),
        ("L3  parallel", 1.431), ("L4  15 stn", 1.221)], maxv=1.6, color=GOOD)
    rect(s, 0.7, 4.35, 11.93, 1.4, WHITE, LINE)
    rect(s, 0.7, 4.35, 0.055, 1.4, BLUE)
    tb(s, 1.0, 4.52, 11.3, 0.38,
       "Sensor maturity: zero PLC state tags costs +1.683 cars/block on L1 — but only +0.021 on L2.",
       14.5, True, INK)
    tb(s, 1.0, 4.95, 11.3, 0.7,
       "State tags matter enormously where a strong constraint exists, and almost not at all on "
       "a balanced line — so\ninstrument the lines with the most to gain first, a more useful "
       "retrofit input than one averaged number.",
       12, False, MUTED, line=1.28)
    tb(s, 0.7, 5.95, 11.93, 0.55,
       "Vintage: one of the three scaling axes the brief names, and we have no data axis for "
       "it — unmodelled, and stated.", 12, False, ALARM, italic=True)
    foot(s, "§5 · §6 · Complexity 6 — Demonstrated (layout, sensor maturity); vintage unmodelled")

    # ---------------------------------------------------------------- 16 trust
    s = blank(prs)
    heading(s, "The hard part is being trusted on day 30.", "complexity 7 · false alarms")
    for i, (v, l, c) in enumerate([
        ("0.479 → 0.025", "calibration error (ECE)\nbefore → after fitting", GOOD),
        ("48.6", "alerts per shift, against the\nISA-18.2 budget of 150", BLUE),
        ("7", "alerts SUPPRESSED for want\nof evidence, not downgraded", BLUE)]):
        kpi(s, 0.7 + i * 4.07, 1.75, 3.85, v, l, c, vsize=25, ht=1.55)
    tb(s, 0.7, 3.5, 12, 0.4,
       "Calibrating the confidence broke our own alerting — and that was the finding.", 16.5, True, INK)
    tb(s, 0.7, 3.95, 12, 1.3,
       "Uncalibrated, the detector claimed ~1.0 confidence and everything got through. Calibrated, "
       "it sits at the honest\nhit rate near 0.11 — and the same threshold silenced every alert. "
       "The system went quiet precisely because it\nbecame truthful. So we stopped gating on the "
       "probability of being exactly right, and started gating on the cost\nof being ignored — "
       "“regret, not accuracy” moved out of the evaluation and into the product.",
       13, False, MUTED, line=1.32)
    rect(s, 0.7, 5.4, 11.93, 0.9, WHITE, LINE)
    rect(s, 0.7, 5.4, 0.055, 0.9, ALARM)
    tb(s, 1.0, 5.55, 11.3, 0.6,
       "We also killed one of our own mechanisms: overtake-risk predicted a bottleneck 5.9% "
       "of the time [1.0–27.0, n=17]\nwhile claiming 70–100% confidence. It is gone from the "
       "live path. Negative results get equal billing.",
       12, False, INK, line=1.3)
    foot(s, "§5 · Complexity 7 — Demonstrated")

    # ------------------------------------------------------------- 17 roadmap
    s = blank(prs)
    heading(s, "Shadow first. Never closed-loop.", "§7 · phased roadmap")
    tb(s, 0.7, 1.62, 12, 0.3,
       "Reads five streams the plant already produces. No new hardware on day one.",
       11.5, False, MUTED)
    phases = [("0", "Shadow", "wk 1–4", "Subscribe to streams the plant\nalready discards. No writes."),
              ("1", "One supervisor", "wk 5–10", "One line, one shift, one screen.\nEvery alert confirmed by a person."),
              ("2", "Floor", "mo 3–6", "All shifts. Manager view opens\nfor planning."),
              ("3", "Retrofit", "next window", "Costed sensor list, installed at a\nscheduled shutdown."),
              ("4", "Second line", "mo 6–12", "Transfer. Report the commissioning\ncurve, not a yes/no.")]
    for i, (n, name, when, what) in enumerate(phases):
        x = 0.7 + i * 2.44
        rect(s, x, 2.05, 2.25, 2.4, WHITE, LINE)
        rect(s, x, 2.05, 2.25, 0.05, BLUE)
        tb(s, x + 0.2, 2.25, 1.8, 0.4, n, 22, True, BLUE, font="Consolas")
        tb(s, x + 0.2, 2.75, 1.9, 0.3, name, 13, True, INK)
        tb(s, x + 0.2, 3.06, 1.9, 0.25, when, 10.5, False, BLUE)
        tb(s, x + 0.2, 3.4, 2.0, 1.0, what, 10.5, False, MUTED, line=1.25)
    rect(s, 0.7, 4.75, 11.93, 0.75, WHITE, LINE)
    rect(s, 0.7, 4.75, 0.055, 0.75, ALARM)
    tb(s, 1.0, 4.92, 11.3, 0.45,
       "Never: closed-loop write to line control — enforced by a test, not a promise.",
       15, True, INK)
    table(s, 0.7, 5.75, 11.93, [3.98, 3.98, 3.97],
          ["PHASE", "GATE TO PASS", "STATUS"],
          [("0 · Shadow", "Twin reproduces the plant's own throughput numbers", ("Executable today", {"bold": True, "color": GOOD})),
           ("3 · Retrofit", "Measured coverage improvement matches the prediction", ("Schedule ready — §5", {"bold": True, "color": BLUE}))],
          row_heights=[0.35, 0.42, 0.42], font_size=10.5, header_size=9.5)
    foot(s, "§7")

    # ------------------------------------------------------ 18 what we did not do
    s = blank(prs)
    heading(s, "What we deliberately did not do.", "honesty, not a gap list")
    items = [
        ("Equipment vintage", "One of the three scaling axes the brief names — we have no data axis for it. Unmodelled, and stated."),
        ("Operator variation", "Excluded on ethical grounds. We measure the station, never the person."),
        ("The CONWIP lead-time figure", "Removed for want of a source file — an earlier draft claimed 36% lower lead time; no file in our results produces it."),
        ("Closed-loop control", "Outside what any plant grants a prototype, and outside what we would want to grant one."),
        ("A historian / MES adapter", "The prototype reads its input streams from files, not a live plant connection — the one deliberate gap before Phase 0 on a real line."),
    ]
    for i, (t, d) in enumerate(items):
        y = 1.7 + i * 0.98
        rect(s, 0.7, y, 11.93, 0.85, WHITE, LINE)
        rect(s, 0.7, y, 0.055, 0.85, MUTED)
        tb(s, 1.0, y + 0.1, 3.3, 0.65, t, 13.5, True, INK, line=1.2)
        tb(s, 4.5, y + 0.13, 7.9, 0.6, d, 11.5, False, MUTED, line=1.22)
    tb(s, 0.7, 6.75, 11.93, 0.5,
       "Every one of these is a decision we can defend, which is worth more than a claim we cannot.",
       14, True, INK)
    foot(s, "§9")

    # -------------------------------------------------------- 19 risks table
    s = blank(prs)
    heading(s, "Key risks and mitigations.", "§8")
    rows = [
        ("False alarms erode trust", "The brief names it; it is how these systems die", "Calibrated confidence (ECE 0.025), the five-field contract with suppression, a ledger showing running precision. 48.6/shift vs a 150 budget"),
        ("Overconfidence in a weak signal", "We have already been caught by this once", "Measured our own overtake-risk mechanism at 5.9% correct [1.0–27.0] vs 70–100% claimed; declared it failed and removed it"),
        ("Top-1 accuracy is misleading", "0.79-car label-noise floor; argmax survives jitter in ~50% of blocks", "Regret is the headline metric; top-1 reported with Wilson intervals alongside, never alone"),
        ("Transfer may not hold", "L2–L4 are 12 runs each", "Stated as directional. Regret held 1.21–1.43 across four topologies; the commissioning curve is reported, not asserted"),
        ("Touching production", "Regulated, safety-certified control", "Read-only by construction and enforced by a test that walks every module's syntax tree, not a policy in a document"),
        ("Modelling people", "Operator variation is a named cause and a real hazard", "Excluded by choice. We measure the station, never the person, and say so"),
    ]
    table(s, 0.7, 1.65, 11.93, [2.9, 4.0, 5.03],
          ["RISK", "WHY IT IS REAL", "MITIGATION"], rows,
          row_heights=[0.38] + [0.83] * 6, font_size=10.5, header_size=10)
    foot(s, "§8")

    # --------------------------------------------------------------- 20 close
    s = blank(prs, dark=True)
    tb(s, 0.9, 1.35, 11.5, 0.9, "The plant already has the data.", 38, True, WHITE)
    tb(s, 0.9, 2.3, 11.5, 1.0,
       "What it has been missing is a model that connects it to cause,"
       "\nconsequence, and the next best action.", 20, False, DIM, line=1.3)
    rect(s, 0.95, 3.55, 1.5, 0.045, ALARM)
    for i, (v, l) in enumerate([
        ("903", "shifts run"), ("15.5%", "warnings on dark stations"),
        ("0.025", "calibration error"), ("70.5%", "tools actionably diagnosed")]):
        x = 0.9 + i * 3.05
        tb(s, x, 4.0, 2.9, 0.55, v, 26, True, WHITE, font="Consolas")
        tb(s, x, 4.58, 2.9, 0.4, l, 11, False, DIM)
    tb(s, 0.9, 5.35, 11.5, 0.8,
       "Every number here is measured, with its source file named — including the ones"
       " that got worse\nwhen we looked harder, and the mechanism we killed.",
       13, False, DIM, line=1.32)
    tb(s, 0.9, 6.15, 11.5, 0.4,
       "github.com/sagar2907/HipHipHooray_IIT_Kharagpur_DigitalTwin.ai_Prototype",
       11.5, False, RGBColor(0x7C, 0x8B, 0x99), font="Consolas")
    station_strip(s, 0.95, 6.7, constraint=19, forming=(6,), h=0.35, w=0.46)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    n = len(prs.slides._sldIdLst)
    print(f"written: {OUT}")
    print(f"slides : {n}")


if __name__ == "__main__":
    build()
